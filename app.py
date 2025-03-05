import streamlit as st
from pptx import Presentation
import os
import google.generativeai as genai
import requests
import json

# Set up Gemini API key
GEMINI_API_KEY = "AIzaSyDumdzIkvU1w8B8ZUmJ2_gyF1cKMw30hH8"
GEMINI_API_URL = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GEMINI_API_KEY}"

def generate_content(topic):
    """Generate slide content using Gemini API"""
    prompt = f"Generate a structured PowerPoint presentation outline for the topic: {topic}. Include slide titles and bullet points."
    
    headers = {"Content-Type": "application/json"}
    data = {
        "contents": [{"parts": [{"text": prompt}]}]
    }
    
    response = requests.post(GEMINI_API_URL, headers=headers, data=json.dumps(data))
    
    if response.status_code == 200:
        response_json = response.json()
        if "candidates" in response_json and len(response_json["candidates"]) > 0:
            return response_json["candidates"][0]["content"]["parts"][0]["text"]
    return "Error generating content"

def create_presentation(slides_data):
    """Generate a PowerPoint presentation based on the AI-generated content"""
    filename = "Smart_Presentation.pptx"
    prs = Presentation()
    
    for slide_info in slides_data.split('\n\n'):
        lines = slide_info.strip().split('\n')
        if len(lines) < 2:
            continue
        
        slide_layout = prs.slide_layouts[1]  # Title and Content Layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = lines[0]  # First line is the title
        content.text = '\n'.join(lines[1:])  # Remaining lines as bullet points
    
    prs.save(filename)
    return filename

# Streamlit UI
st.title("Smart Presentation Generator")
topic = st.text_input("Enter Presentation Topic:")

if st.button("Generate"):
    slides_data = generate_content(topic)
    ppt_filename = create_presentation(slides_data)
    
    st.success("Presentation Generated! Download Below:")
    with open(ppt_filename, "rb") as file:
        st.download_button("Download PPT", file, file_name=ppt_filename)